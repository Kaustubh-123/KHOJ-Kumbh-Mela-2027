"""
KHOJ pitch deck generator -> KHOJ_Pitch_Deck.pptx
Dark, modern theme matching the product UI. Team 21 / Claude Impact Lab.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ---------- palette ----------
BG        = RGBColor(0x0B, 0x11, 0x20)   # near-black slate
BG2       = RGBColor(0x0F, 0x17, 0x2A)   # slate-900
CARD      = RGBColor(0x18, 0x24, 0x3A)   # card fill
CARD_LINE = RGBColor(0x33, 0x41, 0x55)   # slate-700 border
WHITE     = RGBColor(0xF8, 0xFA, 0xFC)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
FAINT     = RGBColor(0x64, 0x74, 0x8B)   # slate-500
RED       = RGBColor(0xEF, 0x44, 0x44)
BLUE      = RGBColor(0x38, 0xBD, 0xF8)
PURPLE    = RGBColor(0xA7, 0x8B, 0xFA)
GREEN     = RGBColor(0x34, 0xD3, 0x99)
AMBER     = RGBColor(0xFB, 0xBF, 0x24)

FONT = "Segoe UI"
FONT_H = "Segoe UI Semibold"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def slide():
    s = prs.slides.add_slide(BLANK)
    # background gradient
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.line.fill.background()
    bg.shadow.inherit = False
    _gradient(bg, BG2, BG, angle=70)
    return s


def _gradient(shape, c1, c2, angle=90):
    sp = shape.fill._xPr
    # remove existing fill
    for tag in ('a:noFill', 'a:solidFill', 'a:gradFill', 'a:blipFill', 'a:pattFill', 'a:grpFill'):
        for el in sp.findall(qn(tag)):
            sp.remove(el)
    grad = sp.makeelement(qn('a:gradFill'), {})
    lst = grad.makeelement(qn('a:gsLst'), {})
    for pos, col in ((0, c1), (100000, c2)):
        gs = grad.makeelement(qn('a:gs'), {'pos': str(pos)})
        clr = grad.makeelement(qn('a:srgbClr'), {'val': '%02X%02X%02X' % (col[0], col[1], col[2])})
        gs.append(clr)
        lst.append(gs)
    grad.append(lst)
    lin = grad.makeelement(qn('a:lin'), {'ang': str(int(angle * 60000)), 'scaled': '1'})
    grad.append(lin)
    # insert before a:ln if present
    ln = sp.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(grad)
    else:
        sp.append(grad)


def rect(s, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
         radius=0.08):
    sp = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, font) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def R(txt, size, color, bold=False, font=FONT):
    return (txt, size, color, bold, font)


def accent_bar(s, x, y, color, w=0.09, h=0.55):
    rect(s, x, y, w, h, fill=color, shape=MSO_SHAPE.RECTANGLE)


def kicker(s, x, y, txt, color):
    text(s, x, y, 8, 0.4, [[R(txt.upper(), 13, color, True)]])


def chip(s, x, y, w, txt, color):
    c = rect(s, x, y, w, 0.42, fill=CARD, line=color, line_w=1.25, radius=0.5)
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_top = Pt(2); tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = txt
    r.font.size = Pt(12); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = FONT
    return c


def footer(s, idx):
    text(s, 0.6, 7.02, 6, 0.3, [[R("KHOJ", 10, FAINT, True), R("  ·  Privacy-Safe AI Triage  ·  Team 21", 10, FAINT)]])
    text(s, 11.2, 7.02, 1.55, 0.3, [[R(f"{idx:02d}", 10, FAINT, True)]], align=PP_ALIGN.RIGHT)


def icon_badge(s, x, y, glyph, color, size=0.72):
    b = rect(s, x, y, size, size, fill=CARD, line=color, line_w=1.5, radius=0.28)
    tf = b.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = glyph
    r.font.size = Pt(22); r.font.color.rgb = color; r.font.bold = True; r.font.name = "Segoe UI Symbol"
    return b


# =========================================================
# SLIDE 1 — TITLE
# =========================================================
s = slide()
# subtle accent glow rectangles
rect(s, -1, -1, 6, 4, fill=None)
text(s, 0.9, 1.55, 11.5, 0.5, [[R("KUMBH MELA 2027  ·  NASHIK", 15, BLUE, True)]])
text(s, 0.85, 2.05, 11.6, 2.0, [[R("KHOJ", 96, WHITE, True, FONT_H)]])
text(s, 0.9, 3.65, 11.5, 0.7, [[R("Privacy-Safe AI Triage Network for Missing Persons", 26, MUTED, False)]])
# tagline rule
rect(s, 0.92, 4.5, 3.2, 0.045, fill=RED, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.9, 4.75, 11.5, 0.6,
     [[R("Turning a missing-person report into a data-driven search plan in under 10 seconds.", 16, WHITE)]])
# chips
chip(s, 0.9, 5.75, 2.7, "Geospatial Triage", BLUE)
chip(s, 3.75, 5.75, 2.5, "Edge-First", PURPLE)
chip(s, 6.4, 5.75, 3.0, "Zero Face Recognition", GREEN)
text(s, 0.9, 6.85, 11.5, 0.4, [[R("Team 21  ·  Claude Impact Lab Hackathon", 13, FAINT, True)]])

# =========================================================
# SLIDE 2 — THE PROBLEM
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, RED)
kicker(s, 0.85, 0.6, "The Challenge", RED)
text(s, 0.82, 0.95, 11.8, 0.9, [[R("Finding one person among ", 34, WHITE, True),
                                  R("tens of millions", 34, RED, True)]])
text(s, 0.85, 1.72, 11.6, 0.6,
     [[R("Kumbh Mela is the largest human gathering on Earth. When someone is lost, every minute compounds the risk — and today's methods have no data behind them.", 15, MUTED)]])

cards = [
    ("⚠", RED, "The Scale", "Shouting in crowds and PA announcements are useless across hundreds of acres and millions of pilgrims."),
    ("\U0001F4DF", AMBER, "The Bottleneck", "Police control rooms get overwhelmed; vague verbal descriptions are slow to act on and easy to lose."),
    ("\U0001F6E1", PURPLE, "The Privacy Trap", "Blanket facial recognition on millions of devotees is an ethical and legal non-starter."),
]
cx = 0.85
for glyph, col, title, body in cards:
    rect(s, cx, 2.7, 3.75, 3.5, fill=CARD, line=CARD_LINE, line_w=1.25, radius=0.06)
    icon_badge(s, cx + 0.4, 3.05, glyph, col)
    text(s, cx + 0.4, 4.0, 3.0, 0.5, [[R(title, 20, WHITE, True, FONT_H)]])
    text(s, cx + 0.4, 4.55, 3.0, 1.5, [[R(body, 13.5, MUTED)]], line_spacing=1.1)
    cx += 4.0
footer(s, 2)

# =========================================================
# SLIDE 3 — THE SOLUTION
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, BLUE)
kicker(s, 0.85, 0.6, "Our Solution", BLUE)
text(s, 0.82, 0.95, 11.8, 0.9, [[R("A distributed network of ", 34, WHITE, True),
                                  R("smart triage kiosks", 34, BLUE, True)]])
text(s, 0.85, 1.72, 11.6, 0.6,
     [[R("Physical terminals at every major ghat and chokepoint. One press, and KHOJ does the rest — on-device, in seconds.", 15, MUTED)]])

steps = [
    ("01", BLUE, "Report in one tap", "A panicking parent or a separated child hits “I am Lost.” The kiosk captures the report instantly — no app, no login."),
    ("02", PURPLE, "Privacy-safe scan", "The AI extracts clothing & colour-pattern metadata locally. No faces, no biometrics are ever stored — only descriptive vectors."),
    ("03", RED, "Geospatial triage", "KHOJ scores every zone by CCTV density, chokepoints & distance, then hands responders an exact, prioritised search plan."),
]
cy = 2.75
for num, col, title, body in steps:
    rect(s, 0.85, cy, 11.6, 1.25, fill=CARD, line=CARD_LINE, line_w=1.0, radius=0.08)
    text(s, 1.1, cy + 0.18, 1.3, 0.9, [[R(num, 40, col, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 2.35, cy + 0.22, 0.04, 0.8, fill=CARD_LINE, shape=MSO_SHAPE.RECTANGLE)
    text(s, 2.7, cy + 0.2, 9.5, 0.5, [[R(title, 21, WHITE, True, FONT_H)]])
    text(s, 2.7, cy + 0.68, 9.5, 0.5, [[R(body, 13.5, MUTED)]])
    cy += 1.4
footer(s, 3)

# =========================================================
# SLIDE 4 — ARCHITECTURE
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, PURPLE)
kicker(s, 0.85, 0.6, "System Architecture", PURPLE)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Edge-first by design", 34, WHITE, True)]])
text(s, 0.85, 1.7, 11.6, 0.5,
     [[R("Everything runs in the browser at the kiosk. No personal data leaves the device — that's the privacy guarantee and the deployability win.", 15, MUTED)]])

flow = [
    ("Report\nForm", BLUE, "Kiosk intake"),
    ("Case\nEngine", PURPLE, "Metadata vectors"),
    ("Recommendation\nEngine", RED, "Zone scoring"),
    ("Live Map +\nDispatch", GREEN, "Action plan"),
]
bx = 0.95
bw = 2.65
for i, (title, col, sub) in enumerate(flow):
    rect(s, bx, 2.75, bw, 1.7, fill=CARD, line=col, line_w=1.5, radius=0.1)
    tb = text(s, bx + 0.15, 2.95, bw - 0.3, 1.0,
              [[R(title, 18, WHITE, True, FONT_H)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, bx + 0.15, 3.95, bw - 0.3, 0.4, [[R(sub, 12, col, True)]], align=PP_ALIGN.CENTER)
    if i < 3:
        text(s, bx + bw - 0.05, 3.05, 0.6, 1.0, [[R("→", 30, FAINT, True)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    bx += bw + 0.5

# data layer
rect(s, 0.95, 4.95, 11.5, 1.5, fill=BG2, line=CARD_LINE, line_w=1.0, radius=0.06)
text(s, 1.25, 5.1, 11, 0.4, [[R("DATA LAYER  —  pre-processed GeoJSON (offline KML → GeoJSON)", 13, MUTED, True)]])
dl = ["zones.geojson", "cctv.geojson", "police.geojson", "chokepoints.geojson", "cases[]  (in-memory)"]
dx = 1.25
for d in dl:
    chip(s, dx, 5.6, 2.05, d, PURPLE)
    dx += 2.18
footer(s, 4)

# =========================================================
# SLIDE 5 — GEOSPATIAL DATA (real numbers)
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, GREEN)
kicker(s, 0.85, 0.6, "Geospatial Intelligence", GREEN)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Built directly on the official datasets", 34, WHITE, True)]])
text(s, 0.85, 1.72, 11.6, 0.5,
     [[R("We parsed the provided KML files into live GeoJSON and plot every point on the command map.", 15, MUTED)]])

stats = [
    ("4,141", "CCTV camera coverage zones", BLUE),
    ("85", "Chokepoints & parking nodes", AMBER),
    ("14", "Police station dispatch points", RED),
]
sx = 0.85
for big, lbl, col in stats:
    rect(s, sx, 2.6, 3.75, 2.2, fill=CARD, line=CARD_LINE, line_w=1.25, radius=0.08)
    rect(s, sx, 2.6, 3.75, 0.12, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, sx + 0.3, 2.95, 3.2, 1.0, [[R(big, 56, col, True, FONT_H)]])
    text(s, sx + 0.32, 4.0, 3.2, 0.6, [[R(lbl, 14, MUTED, True)]])
    sx += 4.0

# pipeline strip
rect(s, 0.85, 5.15, 11.6, 1.3, fill=BG2, line=CARD_LINE, line_w=1.0, radius=0.06)
text(s, 1.15, 5.32, 11, 0.4, [[R("THE PIPELINE", 13, GREEN, True)]])
pipe = [[R("Parse raw .kml  ", 14, WHITE, True), R("→ GeoJSON  ·  ", 14, MUTED),
         R("Haversine ", 14, WHITE, True), R("distance for dynamic search radius  ·  ", 14, MUTED),
         R("Live plot ", 14, WHITE, True), R("on the command dashboard", 14, MUTED)]]
text(s, 1.15, 5.78, 11, 0.5, pipe)
footer(s, 5)

# =========================================================
# SLIDE 6 — THE ALGORITHM
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, AMBER)
kicker(s, 0.85, 0.6, "The Differentiator", AMBER)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Zone Priority Scoring", 34, WHITE, True)]])
text(s, 0.85, 1.7, 11.6, 0.5,
     [[R("Rules-based and fully explainable — not a black box. Every weight has a reason a judge can interrogate.", 15, MUTED)]])

# left: weights table
rect(s, 0.85, 2.55, 6.5, 3.9, fill=CARD, line=CARD_LINE, line_w=1.25, radius=0.06)
text(s, 1.15, 2.75, 6, 0.4, [[R("HOW EACH ZONE IS SCORED", 13, AMBER, True)]])
weights = [
    ("Proximity to last-seen", "1000 / (distance + 1)", BLUE),
    ("CCTV coverage density", "+20 per camera in 300m", GREEN),
    ("Chokepoint concentration", "+15 per chokepoint in 500m", AMBER),
    ("Elderly / child modifier", "+50 if medical point nearby", PURPLE),
    ("Time-since-lost decay", "broadens radius after 60 min", RED),
]
wy = 3.25
for name, formula, col in weights:
    rect(s, 1.15, wy + 0.06, 0.12, 0.36, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, 1.42, wy, 3.4, 0.5, [[R(name, 13.5, WHITE, True)]])
    text(s, 4.6, wy, 2.6, 0.5, [[R(formula, 12, MUTED)]])
    wy += 0.62

# right: output card
rect(s, 7.65, 2.55, 4.8, 3.9, fill=BG2, line=GREEN, line_w=1.5, radius=0.06)
text(s, 7.95, 2.75, 4.2, 0.4, [[R("SAMPLE OUTPUT", 13, GREEN, True)]])
text(s, 7.95, 3.15, 4.2, 0.5, [[R("Elderly male · Zone-C · 30 min", 13, MUTED)]])
out = [
    ("Search radius", "1.5 km", BLUE),
    ("Priority zone", "Zone-C → Gate 7", RED),
    ("Monitor camera", "CCTV-047", GREEN),
    ("Alert station", "Nearest unit", AMBER),
    ("Confidence", "Medium (64%)", PURPLE),
]
oy = 3.7
for k, v, col in out:
    text(s, 7.95, oy, 2.3, 0.4, [[R(k, 13.5, MUTED)]])
    text(s, 9.9, oy, 2.4, 0.4, [[R(v, 14, col, True)]], align=PP_ALIGN.RIGHT)
    rect(s, 7.95, oy + 0.42, 4.2, 0.012, fill=CARD_LINE, shape=MSO_SHAPE.RECTANGLE)
    oy += 0.55
footer(s, 6)

# =========================================================
# SLIDE 7 — KEY FEATURES
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, BLUE)
kicker(s, 0.85, 0.6, "What We Built", BLUE)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Key features", 34, WHITE, True)]])

feats = [
    ("\U0001F198", RED, "“I am Lost” button", "One-tap panic flow for separated pilgrims — scans, captures metadata, and alerts the nearest dispatch instantly."),
    ("\U0001F4F8", BLUE, "AI Fast-Track scan", "Auto-detects clothing vectors from the webcam and simulates the last-known CCTV location."),
    ("\U0001F5FA", PURPLE, "Live Command Center", "Real-time geospatial map plotting every CCTV, police station and chokepoint from the official data."),
    ("\U0001F517", GREEN, "Cross-terminal match", "A report at one kiosk reunites a family at another — the network confirms identity by metadata, not faces."),
]
positions = [(0.85, 2.55), (6.7, 2.55), (0.85, 4.55), (6.7, 4.55)]
for (fx, fy), (glyph, col, title, body) in zip(positions, feats):
    rect(s, fx, fy, 5.6, 1.8, fill=CARD, line=CARD_LINE, line_w=1.0, radius=0.07)
    icon_badge(s, fx + 0.35, fy + 0.35, glyph, col, size=0.85)
    text(s, fx + 1.45, fy + 0.32, 3.9, 0.5, [[R(title, 18, WHITE, True, FONT_H)]])
    text(s, fx + 1.45, fy + 0.82, 3.95, 0.9, [[R(body, 12.5, MUTED)]], line_spacing=1.08)
footer(s, 7)

# =========================================================
# SLIDE 8 — RESPONSIBLE DATA
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, GREEN)
kicker(s, 0.85, 0.6, "Responsible Data", GREEN)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Privacy by ", 34, WHITE, True), R("design", 34, GREEN, True), R(", not by promise", 34, WHITE, True)]])

principles = [
    ("✕", RED, "No face recognition", "Explicitly ruled out. We use descriptive text and geolocation only — never biometrics."),
    ("\U0001F4F6", BLUE, "Footage stays untouched", "We use CCTV camera locations for spatial scoring. No footage is accessed, processed or stored."),
    ("⏱", AMBER, "72-hour retention", "Production data is auto-deleted after the event window. MVP holds nothing past page refresh."),
    ("\U0001F510", PURPLE, "Access-controlled", "Only authenticated police & volunteer coordinators — role-based access from day one."),
]
py = 2.55
for glyph, col, title, body in principles:
    rect(s, 0.85, py, 11.6, 0.95, fill=CARD, line=CARD_LINE, line_w=1.0, radius=0.12)
    icon_badge(s, 1.05, py + 0.16, glyph, col, size=0.62)
    text(s, 1.95, py + 0.14, 3.4, 0.5, [[R(title, 17, WHITE, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 5.4, py + 0.14, 6.8, 0.7, [[R(body, 13, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    py += 1.05
footer(s, 8)

# =========================================================
# SLIDE 9 — LIVE DEMO
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, RED)
kicker(s, 0.85, 0.6, "Live Demo", RED)
text(s, 0.82, 0.95, 11.8, 0.9, [[R("Reunifying a family across the network", 34, WHITE, True)]])
text(s, 0.85, 1.72, 11.6, 0.5, [[R("A single machine simulates a distributed kiosk network via the Global Terminal Selector.", 15, MUTED)]])

acts = [
    ("ACT I", RED, "The Lost Child", "At the Ramkund terminal, the child hits “I am Lost.” KHOJ captures privacy-safe metadata and saves it to the network."),
    ("ACT II", BLUE, "The Parents", "At the Tapovan terminal, parents scan a photo. The screen flashes GREEN — IDENTITY CONFIRMED."),
    ("ACT III", PURPLE, "Command Center", "The global map reveals the child's exact terminal location, ready for dispatch."),
]
ax = 0.85
for tag, col, title, body in acts:
    rect(s, ax, 2.65, 3.75, 3.5, fill=CARD, line=col, line_w=1.5, radius=0.07)
    text(s, ax + 0.35, 2.95, 3.2, 0.4, [[R(tag, 14, col, True)]])
    text(s, ax + 0.35, 3.4, 3.1, 0.6, [[R(title, 21, WHITE, True, FONT_H)]])
    text(s, ax + 0.35, 4.05, 3.1, 1.8, [[R(body, 13.5, MUTED)]], line_spacing=1.12)
    ax += 4.0
footer(s, 9)

# =========================================================
# SLIDE 10 — TECH STACK
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, BLUE)
kicker(s, 0.85, 0.6, "Under the Hood", BLUE)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Tech stack", 34, WHITE, True)]])

stack = [
    ("Kiosk UI", "Next.js 16 · React 19 · Tailwind · Framer Motion", BLUE),
    ("Map & Geo", "React-Leaflet · GeoJSON · @mapbox/togeojson", GREEN),
    ("Camera", "React-Webcam · on-device metadata extraction", PURPLE),
    ("Data store", "localStorage + in-memory store (edge, zero backend)", AMBER),
    ("Prod concept", "FastAPI · PostgreSQL + PostGIS · WebSockets", RED),
]
ty = 2.55
for name, detail, col in stack:
    rect(s, 0.85, ty, 11.6, 0.74, fill=CARD, line=CARD_LINE, line_w=1.0, radius=0.14)
    rect(s, 0.85, ty, 0.14, 0.74, fill=col, shape=MSO_SHAPE.RECTANGLE)
    text(s, 1.25, ty + 0.1, 3.2, 0.5, [[R(name, 16, WHITE, True, FONT_H)]], anchor=MSO_ANCHOR.MIDDLE)
    text(s, 4.5, ty + 0.1, 7.7, 0.5, [[R(detail, 14, MUTED)]], anchor=MSO_ANCHOR.MIDDLE)
    ty += 0.84
footer(s, 10)

# =========================================================
# SLIDE 11 — ROADMAP
# =========================================================
s = slide()
accent_bar(s, 0.6, 0.62, PURPLE)
kicker(s, 0.85, 0.6, "Beyond 2027", PURPLE)
text(s, 0.82, 0.95, 11.8, 0.8, [[R("Where KHOJ goes next", 34, WHITE, True)]])

phases = [
    ("NOW", GREEN, "MVP", "Geospatial triage on real data, running live at the edge — what you'll see today."),
    ("NEXT", BLUE, "Vernacular Voice AI", "Voice-to-text in local dialects for illiterate or panicking pilgrims."),
    ("LATER", PURPLE, "Drone Dispatch", "Ping kiosk coordinates directly to automated aerial units for visual confirmation."),
    ("SCALE", RED, "Sensor Fusion", "Live gate-sensor crowd density + volunteer coordination via messaging APIs."),
]
px = 0.85
# timeline line
rect(s, 1.0, 3.05, 11.3, 0.03, fill=CARD_LINE, shape=MSO_SHAPE.RECTANGLE)
for tag, col, title, body in phases:
    rect(s, px + 1.25, 2.95, 0.22, 0.22, fill=col, shape=MSO_SHAPE.OVAL)
    rect(s, px, 3.5, 2.85, 2.7, fill=CARD, line=CARD_LINE, line_w=1.0, radius=0.07)
    text(s, px + 0.3, 3.75, 2.3, 0.4, [[R(tag, 13, col, True)]])
    text(s, px + 0.3, 4.2, 2.3, 0.8, [[R(title, 18, WHITE, True, FONT_H)]])
    text(s, px + 0.3, 5.05, 2.3, 1.1, [[R(body, 12.5, MUTED)]], line_spacing=1.1)
    px += 2.95
footer(s, 11)

# =========================================================
# SLIDE 12 — THANK YOU
# =========================================================
s = slide()
text(s, 0.9, 2.45, 11.5, 0.5, [[R("KUMBH MELA 2027  ·  NASHIK", 14, BLUE, True)]])
text(s, 0.85, 2.85, 11.6, 1.6, [[R("Thank you", 76, WHITE, True, FONT_H)]])
rect(s, 0.92, 4.35, 3.2, 0.045, fill=RED, shape=MSO_SHAPE.RECTANGLE)
text(s, 0.9, 4.6, 11.5, 0.6,
     [[R("Because every second in a crowd of a million people matters.", 18, MUTED)]])
text(s, 0.9, 5.6, 11.5, 0.4, [[R("Team 21  ·  Claude Impact Lab Hackathon", 14, FAINT, True)]])
chip(s, 0.9, 6.1, 2.2, "Questions?", GREEN)

prs.save("KHOJ_Pitch_Deck.pptx")
print("Saved KHOJ_Pitch_Deck.pptx with", len(prs.slides._sldIdLst), "slides")
